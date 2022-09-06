import json
import os
import subprocess
import tempfile
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor
from glob import glob
from io import BytesIO
from multiprocessing import Pool, cpu_count
from os import mkdir, path, rename, rmdir
from posixpath import splitext
from re import I
from shutil import copy, rmtree
from tempfile import gettempdir
from threading import Thread
from zipfile import ZIP_DEFLATED, BadZipFile, ZipFile

import aspose.slides as slides
from alive_progress import alive_bar
from PIL import Image, ImageFile, ImageOps
from pptx import opc
from pptx.api import Presentation  # type: ignore
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR  # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.parts import image
from pptx.presentation import Presentation as PPTX
from pptx.util import Inches

# opc.package.Part.related_parts = related_parts

# TEMP_DIR = gettempdir()
ImageFile.LOAD_TRUNCATED_IMAGES = True


def open_file(filename: str):
    '''Open document with default application in Python.'''
    try:
        os.startfile(filename)  # type: ignore
    except AttributeError:
        subprocess.call(['open', filename])


def file_list(location: str):
    return glob(location, recursive=True)


def create_dir(path: str):
    try:
        mkdir(path)
    except:
        pass


# FILE = path.join("input", "powerpoints",
    #  "'k Ben een koninklijk kind (EL 180).pptx")

# files = file_list(path.join("input"))

# # 0. Clear all previous operations
# in_zips_path = path.join(TEMP_DIR, "in-zips")
# extracted_path = path.join(TEMP_DIR, "extracted")
# if path.exists(in_zips_path):
#     rmtree(in_zips_path)
# if path.exists(extracted_path):
#     rmtree(extracted_path)

def clear_output():
    rmtree("images")
    mkdir('images')
    rmtree("output")
    mkdir("output")
    rmtree("between")
    mkdir("between")

# 1. Convert powerpoint to zip


in_zips_path = tempfile.mkdtemp()
extracted_path = tempfile.mkdtemp()


def extract_media(file: str):

    # create_dir(extracted_path)

    # Create a directory for the all the media
    file_name = path.split(path.splitext(file)[0])[1]
    new_dir = path.join(extracted_path, file_name)
    create_dir(new_dir)

    # Read the zipfile
    with ZipFile(file, "r") as zip:
        # Get a list of all files
        for info in zip.infolist():
            # Check for a media directory
            if("media" in info.filename):
                # Extract all files to a new path
                new_path = zip.extract(
                    info.filename, extracted_path)
                # Move the files
                move(new_path, new_dir)

    rmtree(path.join(extracted_path, "ppt"))
    return new_dir


# extract_media(FILE)
# 3. Remove all white pixels from all media AND invert image


def remove_white(image):
    image = image.convert("RGBA")
    data = image.getdata()

    newData = []
    for item in data:
        if item[0] == 255 and item[1] == 255 and item[2] == 255:
            newData.append((255, 255, 255, 0))
        else:
            newData.append(item)

    image.putdata(newData)
    return image


def invert_image(image):
    # image = Image.open(image_file)
    image = image.convert("RGBA")

    r, g, b, a = image.split()
    rgb_image = Image.merge('RGB', (r, g, b))

    inverted_image = ImageOps.invert(rgb_image)

    r2, g2, b2 = inverted_image.split()

    image = Image.merge('RGBA', (r2, g2, b2, a))
    return image


def remove_transparency(image, bg_colour=(255, 255, 255)):
    if image.mode in ('RGBA', 'LA') or (image.mode == 'P' and 'transparency' in image.info):
        alpha = image.convert('RGBA').split()[-1]
        bg = Image.new("RGBA", image.size, bg_colour + (255,))
        bg.paste(image, mask=alpha)
        return bg
    else:
        return image


def replace_media(folder: str):
    file = f"{folder}.zip"

    input_path = path.join(in_zips_path, file)

    # out_folder = path.join(extracted_path, "out_folder")

    # create_dir(out_folder)
    output_path = path.join(extracted_path, file)

    with ZipFile(input_path, "r") as in_zip:
        with ZipFile(output_path, "w") as out_zip:
            out_zip.comment = in_zip.comment
            for info in in_zip.infolist():
                if("media" in info.filename):
                    location = path.split(info.filename)[1]
                    new_location = path.join(
                        extracted_path, folder, location)

                    out_zip.write(
                        new_location, info.filename)
                else:
                    out_zip.writestr(
                        info, in_zip.read(info.filename))
    return output_path


# FOLDER = path.splitext(path.split(FILE)[1])[0]

# out_zip = replace_media(FOLDER)

# 6. Convert the zip folders back to powerpoints

def change_extension(file: str, ext: str):
    # Change the extension of a file and rename it
    new_file = f"{path.splitext(file)[0]}.{ext}"
    rename(file, new_file)
    return new_file


def convert_to_pptx(file: str):
    dir = path.join("output")
    create_dir(dir)
    return change_extension(new_file, "pptx")


# Replace the current image on the slide with the new inverted_image


# pptx_file = convert_to_pptx(out_zip)


# 7. Update the powerpoint colors


# function to replace a picture in a slide with a new picture
def replace_picture(slide, shape, new_picture):
    """
    function to replace a picture in a slide with a new picture
    """
    if(shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER):

        # get the placeholder shape
        placeholder = slide.shapes[shape.placeholder_format.idx]
        # get the placeholder size
        width, height = placeholder.width, placeholder.height
        # get the placeholder position
        left, top = placeholder.left, placeholder.top
        # delete the placeholder
        placeholder.delete()
        # add the new picture
        slide.shapes.add_picture(new_picture, left, top,
                                 width=width, height=height)

# def replace_picture(slide_obj, shape_obj, img_location):

    # with open(img_location, "rb") as file_obj:
    #     r_img_blob = file_obj.read()
    #     img_pic = shape_obj._pic
    #     img_rid = img_pic.xpath("./p:blipFill/a:blip/@r:embed")[0]
    #     img_part = slide_obj.part.related_parts[img_rid]
    #     img_part._blob = r_img_blob
    # noinspection PyProtectedMember
    # new_shape = slide_obj.shapes.add_picture(
    #     img_location,
    #     Inches(shape_obj.left),
    #     Inches(shape_obj.top),
    #     Inches(shape_obj.width),
    #     Inches(shape_obj.height),
    # )
    # old_pic = shape_obj._element
    # new_pic = new_shape._element
    # old_pic.addnext(new_pic)
    # old_pic.getparent().remove(old_pic)


def extract_images(name: str, pptx: PPTX):

    # name: str = pptx.core_properties.subject

    images_folder_path = path.join("images", name)
    rmtree(images_folder_path, ignore_errors=True)
    create_dir(images_folder_path)

    for index, slide in enumerate(pptx.slides):  # type: ignore
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Change the background fill of the picture to white
                # slide_part, rId = shape.part, shape._element.blip_rId
                # print(slide_part)
                # image_part = slide_part.related_parts[rId]
                # image_part.blob = new_blob
                image = shape.image
                image_bytes = image.blob

                image_filename = f'{name}-{index}.{image.ext}'

                image_path = path.join(tempfile.mkdtemp(), image_filename)

                with open(image_path, 'wb') as f:
                    f.write(image_bytes)

                open_image = Image.open(image_path)

                transparent_image = remove_white(open_image)
                inverted_image = invert_image(transparent_image)

                new_image_path = path.join(
                    images_folder_path, f"image{index + 1}.png")
                inverted_image.save(new_image_path, "PNG")

    return pptx, images_folder_path


def invert_colors(pptx: PPTX):
    for slide in pptx.slides:  # type: ignore
        try:
            # Update background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.patterned()
            fill.back_color.rgb = RGBColor(0, 0, 0)

            # for shape in slide.placeholders:
            #     print('%d %s' % (shape.placeholder_format.idx, shape.name))

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)

        except KeyError:
            continue

    return pptx

# modify_pptx(pptx_file)

# print(TEMP_DIR)


def convert_image(file: str):
    image = Image.open(file)

    image = image.convert("RGBA")

    r, g, b, a = image.split()
    rgb_image = Image.merge('RGB', (r, g, b))

    inverted_image = ImageOps.invert(rgb_image)

    r2, g2, b2 = inverted_image.split()

    image = Image.merge('RGBA', (r2, g2, b2, a))

    image.save(file, "PNG")


def convert_pptx(file: str):

    name = path.split(path.splitext(file)[0])[1]
    between_path = path.join("between", f"{name}.pptx")
    output_path = path.join("output", f"{name}.pptx")
    create_dir(path.join("images", name))
    # with slides.Presentation(file) as presentation:
    #     for idx, image in enumerate(presentation.images):
    #         opened_image = Image.open(BytesIO(image.binary_data))

    #         # removed_white_image = remove_white(opened_image)
    #         inverted_image = invert_image(opened_image)

    #         image_path = path.join(
    #             "images", name, f"image{idx + 1}.png")

    #         inverted_image.save(image_path, "PNG")

    #         image.replace_image(read_all_bytes(image_path))

    #     presentation.save(
    #         between_path, slides.export.SaveFormat.PPTX)
    pptx: PPTX = Presentation(file)

    # pptx: PPTX = Presentation(between_path)

    pptx, images_folder_path = extract_images(name, pptx)
    pptx = invert_colors(pptx)

    pptx.save(between_path)

    input_zip_file = convert_to_zip(between_path)

    output_zip_file = path.join(
        "output", path.split(input_zip_file)[1])

    copy(input_zip_file, output_zip_file)

    with ZipFile(input_zip_file, "r") as zip_ref:
        with ZipFile(output_zip_file, "w") as zip:
            zip.comment = zip_ref.comment
            for info in zip_ref.infolist():
                # replace the old images with the new images
                if("media" in info.filename):
                    image_name = path.split(info.filename)[1]
                    image_path = path.join(
                        images_folder_path, image_name)

                    zip.write(image_path, info.filename)
                else:
                    zip.writestr(
                        info, zip_ref.read(info.filename))

    change_extension(output_zip_file, "pptx")

    return file


# def read_image(image_file: str, gray_scale=False):
#     image_src = cv2.imread(image_file)
#     if gray_scale:
#         image_src = cv2.cvtColor(image_src, cv2.COLOR_BGR2GRAY)
#     else:
#         image_src = cv2.cvtColor(image_src, cv2.COLOR_BGR2RGB)
#     return image_src


def convert_to_zip(file: str):
    return change_extension(file, "zip")


def read_all_bytes(file_name):
    with open(file_name, "rb") as stream:
        return stream.read()


if __name__ == "__main__":
    clear_output()

    files = file_list("input/**/*.pptx")

    with alive_bar(len(files), dual_line=True,  title="Converting powerpoints...") as bar:
        with Pool() as pool:
            for file in pool.imap_unordered(convert_pptx, files):
                bar.text = f"Converted '{file}'..."
                bar()
        # for file in files:

    # file = file_list(path.join("output", "**/*.pptx"))[0]

    # open_file(file)
