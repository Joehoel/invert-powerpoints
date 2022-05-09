from glob import glob
from tempfile import gettempdir
from os import rename, path, mkdir, rmdir
from shutil import copy, move, rmtree
from PIL import Image, ImageOps
from zipfile import BadZipFile, ZipFile
from tqdm import tqdm
from pptx.api import Presentation
from pptx.presentation import Presentation as PPTX
from pptx.dml.color import RGBColor


TEMP_DIR = gettempdir()


def file_list(location: str):
    return glob(location, recursive=True)


def create_dir(path: str):
    try:
        mkdir(path)
    except:
        pass


def change_extension(location: str, ext: str):
    base = path.splitext(location)[0]
    file = f"{base}.{ext}"
    rename(location, file)
    return file


FILE = path.join("input", "powerpoints",
                 "'k Ben een koninklijk kind (EL 180).pptx")

# 0. Clear all previous operations
rmtree("output")
mkdir("output")

# 1. Convert powerpoint to zip


def convert_to_zip(file: str):
    dir = path.join(TEMP_DIR, "in-zips")
    create_dir(dir)
    new_file = copy(file, dir)
    change_extension(new_file, "zip")


convert_to_zip(FILE)

# 2. Extract media from zip


def extract_media(file: str):
    # # print(f"👽 Extracting media from {location}...")

    extracted_path = path.join("output", "extracted")
    create_dir(extracted_path)

    # Create a directory for the all the media
    file_name = path.split(path.splitext(file)[0])[1]
    new_dir = path.join(extracted_path, file_name)
    create_dir(new_dir)

    # Read the zipfile
    with ZipFile(file, "r") as zip:
        # Get a list of all files
        for info in zip.infolist():
            # Check for a media directory
            if("media" in info.filename):
                # Extract all files to a new path
                new_path = zip.extract(
                    info.filename, extracted_path)
                # Move the files
                move(new_path, new_dir)

    rmtree(path.join(extracted_path, "ppt"))


extract_media(FILE)
# 3. Remove all white pixels from all media AND invert image


def remove_white(file: str):
    image = Image.open(file)

    image = image.convert("RGBA")
    data = image.getdata()

    newData = []
    for item in data:
        if item[0] == 255 and item[1] == 255 and item[2] == 255:
            newData.append((255, 255, 255, 0))
        else:
            newData.append(item)

    image.putdata(newData)
    return image.convert("RGB")


def invert_image(file: str):
    image = Image.open(file)
    image = image.convert("RGBA")

    r, g, b, a = image.split()
    rgb_image = Image.merge('RGB', (r, g, b))

    inverted_image = ImageOps.invert(rgb_image)

    r2, g2, b2 = inverted_image.split()

    image = Image.merge('RGBA', (r2, g2, b2, a))
    return image


media = glob(path.join("output", "extracted", "**/*"))

print("Removing all white pixels from all media and inverting the image")
for file in tqdm(media):
    image = remove_white(file)
    image = invert_image(file)
    image.save(file)


# 5. Replace the new images with the old images in the zip


def replace_media(folder: str):
    file = f"{folder}.zip"

    input_path = path.join(TEMP_DIR, "in-zips", file)

    out_folder = path.join(TEMP_DIR, "in_zips")

    create_dir(out_folder)
    output_path = path.join(out_folder, file)

    with ZipFile(input_path, "r") as in_zip:
        with ZipFile(output_path, "w") as out_zip:
            out_zip.comment = in_zip.comment
            for info in in_zip.infolist():
                if("media" in info.filename):
                    location = path.split(info.filename)[1]
                    new_location = path.join(
                        "output", "extracted", folder, location)

                    out_zip.write(
                        new_location, info.filename)
                else:
                    out_zip.writestr(
                        info, in_zip.read(info.filename))
    return output_path


FOLDER = path.splitext(path.split(FILE)[1])[0]

out_zip = replace_media(FOLDER)

# 6. Convert the zip folders back to powerpoints


def convert_to_pptx(file: str):
    dir = path.join("output", "powerpoints")
    create_dir(dir)
    new_file = copy(file, dir)
    return change_extension(new_file, "pptx")


pptx_file = convert_to_pptx(out_zip)

# 7. Update the powerpoint colors


def modify_pptx(file: str):
    # print(f"✒️ Modifying: ./output/powerpoints/{file}.pptx")
    pptx: PPTX = Presentation(file)

    for slide in pptx.slides:
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        run = p.add_run()
                        run.font.color.rgb = RGBColor(255, 255, 255)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.patterned()
            fill.back_color.rgb = RGBColor(0, 0, 0)

        except KeyError:
            # print(f"🗑️ Corrupt file: {file}.pptx")
            continue

    pptx.save(file)


modify_pptx(pptx_file)
